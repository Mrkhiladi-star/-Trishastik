import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Theme Palette (Trishastik Dark Premium style)
    bg_color = RGBColor(9, 15, 29)          # Deep Slate-950 (Dark background)
    emerald = RGBColor(16, 185, 129)       # Emerald-500 (Primary accent)
    light_emerald = RGBColor(52, 211, 153)  # Emerald-400 (Secondary highlight)
    white = RGBColor(255, 255, 255)         # Pure White (Headers and emphasis)
    slate_300 = RGBColor(226, 232, 240)     # Light Gray (Body text)
    slate_500 = RGBColor(100, 116, 139)     # Muted Gray (Footers)
    card_bg = RGBColor(21, 32, 54)          # Deep blue-gray (Card background)
    card_border = RGBColor(38, 55, 87)      # Muted slate (Card border)
    gold = RGBColor(245, 158, 11)           # Amber-500 (Alerts & key values)
    
    def add_blank_slide(bg_col=bg_color):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg_col
        rect.line.color.rgb = bg_col
        return slide

    def add_header(slide, title_text):
        # Title Box - starting at 0.35, ending at 1.05
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.35), Inches(8.5), Inches(0.7))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Arial'
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = white
        
        # Muted Context Header
        sub_box = slide.shapes.add_textbox(Inches(9.5), Inches(0.35), Inches(3.08), Inches(0.4))
        sub_tf = sub_box.text_frame
        sub_tf.word_wrap = True
        sub_tf.margin_left = sub_tf.margin_top = sub_tf.margin_right = sub_tf.margin_bottom = 0
        sub_p = sub_tf.paragraphs[0]
        sub_p.text = "Trishastik Bharat | IIT BHU"
        sub_p.font.name = 'Arial'
        sub_p.font.size = Pt(11)
        sub_p.font.bold = True
        sub_p.font.color.rgb = emerald
        sub_p.alignment = PP_ALIGN.RIGHT

        # Accent Line - Placed at 1.35 to prevent any overlap with the title text
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.35), Inches(11.83), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = emerald
        line.line.color.rgb = emerald
        
        # Footer
        footer_box = slide.shapes.add_textbox(Inches(0.75), Inches(7.0), Inches(11.83), Inches(0.3))
        ftf = footer_box.text_frame
        ftf.word_wrap = True
        fp = ftf.paragraphs[0]
        fp.text = "Presentation on Trishastik Bharat AgriTech Platform | Solving Core Agricultural Inefficiencies"
        fp.font.name = 'Arial'
        fp.font.size = Pt(9)
        fp.font.color.rgb = slate_500

    def add_card(slide, left, top, width, height, title, body_bullets, title_color=emerald):
        # Card Background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = card_bg
        card.line.color.rgb = card_border
        card.line.width = Pt(1.5)
        
        # Text Frame inside card
        tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), width - Inches(0.5), height - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        # Card Title
        if title:
            tp = tf.paragraphs[0]
            tp.text = title
            tp.font.name = 'Arial'
            tp.font.size = Pt(16)
            tp.font.bold = True
            tp.font.color.rgb = title_color
            tp.space_after = Pt(10)
            
        # Bullet points
        for i, bullet in enumerate(body_bullets):
            p = tf.add_paragraph() if (title or i > 0) else tf.paragraphs[0]
            p.text = "• " + bullet
            p.font.name = 'Arial'
            p.font.size = Pt(11.5)
            p.font.color.rgb = slate_300
            p.space_after = Pt(6)

    def add_diagram_slide(slide_title, steps_data):
        # Helper to create standardized 4-step horizontal flowchart diagrams
        slide = add_blank_slide()
        add_header(slide, slide_title)
        
        box_width = Inches(2.55)
        box_height = Inches(4.2)
        box_top = Inches(2.0)
        gap = Inches(0.5)
        
        for idx, (step_header, bullets, color) in enumerate(steps_data):
            box_left = Inches(0.75) + idx * (box_width + gap)
            
            # Step Card Shape
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_left, box_top, box_width, box_height)
            card.fill.solid()
            card.fill.fore_color.rgb = card_bg
            card.line.color.rgb = card_border
            card.line.width = Pt(1.5)
            
            # Number Badge inside Card
            num_box = slide.shapes.add_textbox(box_left + Inches(0.2), box_top + Inches(0.15), Inches(0.6), Inches(0.4))
            ntf = num_box.text_frame
            ntf.margin_left = ntf.margin_top = ntf.margin_right = ntf.margin_bottom = 0
            np = ntf.paragraphs[0]
            np.text = f"0{idx+1}"
            np.font.name = 'Arial'
            np.font.size = Pt(16)
            np.font.bold = True
            np.font.color.rgb = color
            
            # Bullet/Text Box
            text_box = slide.shapes.add_textbox(box_left + Inches(0.2), box_top + Inches(0.65), box_width - Inches(0.4), box_height - Inches(0.8))
            tf = text_box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
            
            # Step Title
            tp = tf.paragraphs[0]
            tp.text = step_header
            tp.font.name = 'Arial'
            tp.font.size = Pt(13)
            tp.font.bold = True
            tp.font.color.rgb = white
            tp.space_after = Pt(10)
            
            # Description bullets
            for b_idx, bullet in enumerate(bullets):
                bp = tf.add_paragraph()
                bp.text = "• " + bullet
                bp.font.name = 'Arial'
                bp.font.size = Pt(9.5)
                bp.font.color.rgb = slate_300
                bp.space_after = Pt(4)
                
            # Right Connection Arrow
            if idx < 3:
                arrow_left = box_left + box_width + Inches(0.05)
                arrow_top = box_top + Inches(1.9)
                arrow_width = gap - Inches(0.1)
                
                arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, arrow_top, arrow_width, Inches(0.3))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = color
                arrow.line.color.rgb = color
                
        return slide

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Cover)
    # -------------------------------------------------------------
    slide1 = add_blank_slide()
    
    # Decorative visual background glow
    glow = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.0), Inches(-1.5), Inches(6.0), Inches(6.0))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(6, 78, 59)
    glow.line.fill.background()
    
    # Title Text Box
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "TRISHASTIK BHARAT"
    p1.font.name = 'Georgia'
    p1.font.size = Pt(54)
    p1.font.bold = True
    p1.font.color.rgb = emerald
    p1.space_after = Pt(8)
    
    p2 = tf.add_paragraph()
    p2.text = "Sustainable Direct AgriTech & Circular Farming Ecosystem"
    p2.font.name = 'Arial'
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = white
    p2.space_after = Pt(18)
    
    p3 = tf.add_paragraph()
    p3.text = "Empowering Farmers through Disintermediated Organic Markets, Shared Rentals,\nDirect-to-Mill Inputs, Monitored Soil Testing, and Auto-Routed Logistics."
    p3.font.name = 'Arial'
    p3.font.size = Pt(14)
    p3.font.color.rgb = slate_300
    
    # Decorative Line
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.7), Inches(5.0), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = emerald
    line.line.color.rgb = emerald

    # Presenter credits
    pres_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(10.0), Inches(1.8))
    ptf = pres_box.text_frame
    ptf.word_wrap = True
    pp = ptf.paragraphs[0]
    pp.text = "IIT BHU (Varanasi) | Project Presentation"
    pp.font.name = 'Arial'
    pp.font.size = Pt(12)
    pp.font.bold = True
    pp.font.color.rgb = white
    pp.space_after = Pt(4)
    
    pp2 = ptf.add_paragraph()
    pp2.text = "Designed to Optimize Agricultural Value Chains & Eliminate Middlemen"
    pp2.font.name = 'Arial'
    pp2.font.size = Pt(11)
    pp2.font.color.rgb = light_emerald

    # -------------------------------------------------------------
    # SLIDE 2: Core Vision & Problem Statement
    # -------------------------------------------------------------
    slide2 = add_blank_slide()
    add_header(slide2, "The Agrarian Crisis & Trishastik Vision")
    
    # Left column text (The Problem) - Placed at 1.85
    left_box = slide2.shapes.add_textbox(Inches(0.75), Inches(1.85), Inches(5.6), Inches(4.8))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    
    lp1 = ltf.paragraphs[0]
    lp1.text = "Core Structural Inefficiencies"
    lp1.font.name = 'Arial'
    lp1.font.size = Pt(20)
    lp1.font.bold = True
    lp1.font.color.rgb = gold
    lp1.space_after = Pt(12)
    
    problems = [
        "Exploitative Middlemen: Intermediaries capture massive margins, leaving farmers with low payouts and high debt.",
        "Chemical Monopoly: Online agricultural stores focus exclusively on chemicals; direct organic inputs are hard to procure.",
        "Unmonitored Testing: Free government labs exist but have zero operational monitoring or tracking, creating long delays.",
        "Prohibitive CapEx: Smallholder farmers cannot afford to purchase tractors and heavy machinery, limiting mechanization."
    ]
    for prob in problems:
        p = ltf.add_paragraph()
        p.text = "• " + prob
        p.font.name = 'Arial'
        p.font.size = Pt(12)
        p.font.color.rgb = slate_300
        p.space_after = Pt(10)
        
    # Right column card (The Solution/Vision)
    add_card(
        slide2, 
        Inches(6.98), Inches(1.85), Inches(5.6), Inches(4.8),
        "Trishastik Bharat: The Solution",
        [
            "Direct-to-Buyer Market: Bypasses intermediate brokers to trade organic produce directly at fair market rates.",
            "Mill-to-Farm Sourcing: Direct online access to organic waste byproducts (press mud, compost) from sugar and paper mills.",
            "Monitored Soil Health: Portal tracks soil agents and city labs in real-time, enforcing accountability for fast testing.",
            "P2P Equipment Rentals: Access to shared machinery on demand, turning major capital costs into cheap operating expenses."
        ],
        title_color=white
    )

    # -------------------------------------------------------------
    # SLIDE 3: Direct Farmer-to-Buyer Organic Market (No Middlemen)
    # -------------------------------------------------------------
    slide3 = add_blank_slide()
    add_header(slide3, "Direct Organic Trading: Bypassing the Middlemen")
    
    # Left Card: Traditional Market Exploitation
    add_card(
        slide3,
        Inches(0.75), Inches(1.85), Inches(5.6), Inches(4.8),
        "Traditional Route (Exploitative)",
        [
            "High Broker Cuts: Middlemen capture 25%+ of crop market value.",
            "Concrete Example (Wheat Pricing):\n  - Government MSP Rate:  ₹2,685 / Quintal\n  - Middleman Buying Rate: ₹2,000 / Quintal\n  - Direct Loss to Farmer:  ₹685 / Quintal (Over 25% stolen!)",
            "Payment Delays: Farmers wait weeks for cash, leading to debt.",
            "Lack of Transparency: Weight manipulation and zero price control."
        ],
        title_color=gold
    )
    
    # Right Card: Trishastik Direct Marketplace
    add_card(
        slide3,
        Inches(6.98), Inches(1.85), Inches(5.6), Inches(4.8),
        "Trishastik Direct Route (Fair & Transparent)",
        [
            "Direct-to-Buyer Contracts: Connects farmers straight to certified organic buyers, food companies, and retailers.",
            "Full Payout Preservation:\n  - Farmer gets the actual market value (e.g. ₹2,685+ for wheat).\n  - Escrow payments hold buyer funds and release them instantly.",
            "No Commission Pockets: Zero middleman margins. 100% of fair value reaches the agricultural producer.",
            "Quality Premium: Farmers earn higher premiums for verified organic crops certified via our portal."
        ],
        title_color=emerald
    )

    # -------------------------------------------------------------
    # SLIDE 4: Equipment Rental & Portal Purchase
    # -------------------------------------------------------------
    slide4 = add_blank_slide()
    add_header(slide4, "Equipment Portal: Pay-Per-Use Rentals & Buying")
    
    # Left Column: Rental Model
    add_card(
        slide4,
        Inches(0.75), Inches(1.85), Inches(5.6), Inches(4.8),
        "1. Peer-to-Peer Rental Network",
        [
            "Affordable Mechanization: Farmers rent tractors, seeders, and harvesters by the hour rather than buying them.",
            "Up to 80% Cost Reduction: Converts crushing capital expenditures (CapEx) into flexible, low operational costs (OpEx).",
            "Monetize Idle Machinery: Farmers owning equipment can list it for rent, generating secondary income streams.",
            "Geospatial Search: Finds the nearest available equipment using Leaflet-mapped farm listings."
        ],
        title_color=emerald
    )
    
    # Right Column: Purchase Model
    add_card(
        slide4,
        Inches(6.98), Inches(1.85), Inches(5.6), Inches(4.8),
        "2. Direct Portal Equipment Sales",
        [
            "Zero Dealer Markups: Partners directly with manufacturing brands to sell new machinery and tools online.",
            "Transparent Factory Prices: Cuts out heavy regional distributor fees and retail commissions.",
            "Integrated Quality Assurance: Only certified, high-performance machinery is permitted to list.",
            "Micro-Financing Access: Integrates cooperative bank loans directly on the purchase checkout screen."
        ],
        title_color=light_emerald
    )

    # -------------------------------------------------------------
    # SLIDE 5: Direct Mill-to-Farm Organic Fertilizers
    # -------------------------------------------------------------
    fertilizer_steps = [
        (
            "1. Sourcing from Mills",
            [
                "Partner directly with local sugar, paper, and food mills.",
                "Source nutrient-rich byproducts like press mud and spent wash compost."
            ],
            emerald
        ),
        (
            "2. Portal Integration",
            [
                "Mills list organic inputs in bulk directly on Trishastik.",
                "Solves the online chemical monopoly; lists factory-certified organic compost."
            ],
            light_emerald
        ),
        (
            "3. Direct Purchase",
            [
                "Farmers buy fertilizers directly from mills at bulk factory rates.",
                "Bypasses standard wholesale fertilizer dealers."
            ],
            white
        ),
        (
            "4. Farm Delivery",
            [
                "Auto-dispatched logistics deliver bulk orders straight to the farm grid.",
                "Lowers overall input costs while restoring soil health."
            ],
            emerald
        )
    ]
    add_diagram_slide("Mill-to-Farm Organic Fertilizers: Sourced Direct from Mills", fertilizer_steps)

    # -------------------------------------------------------------
    # SLIDE 6: Accountable & Monitored Soil Testing
    # -------------------------------------------------------------
    soil_steps = [
        (
            "1. Portal Request",
            [
                "Farmer requests test, pins GPS coordinates on Leaflet Map.",
                "Saves planned crop targets and history to user profile."
            ],
            emerald
        ),
        (
            "2. Agent Dispatch",
            [
                "System tracks assigned field agents for that area.",
                "Audits agent response time to ensure samples are collected."
            ],
            light_emerald
        ),
        (
            "3. Lab Test Upload",
            [
                "Free government city labs test soil chemical structure.",
                "Parsed chemical metrics and raw PDF are uploaded to portal."
            ],
            white
        ),
        (
            "4. Grok AI Diagnostics",
            [
                "AI engine evaluates N-P-K and pH levels.",
                "Pushes customized organic fertilizer dosage to farmer dashboard."
            ],
            emerald
        )
    ]
    add_diagram_slide("Monitored Soil Diagnostics: Bridging the Accountability Gap", soil_steps)

    # -------------------------------------------------------------
    # SLIDE 7: Farmers as Marketplace Sellers (Direct to Consumer)
    # -------------------------------------------------------------
    slide7 = add_blank_slide()
    add_header(slide7, "Marketplace Model: Farmers as Retail Sellers")
    
    # 3 columns setup
    col_width = Inches(3.6)
    col_gap = Inches(0.5)
    col_top = Inches(1.85)
    col_height = Inches(4.8)
    
    add_card(
        slide7,
        Inches(0.75) + 0 * (col_width + col_gap), col_top, col_width, col_height,
        "1. Storefront Creation",
        [
            "Easy Onboarding: Farmers switch their dashboard profile to 'Seller' in one click.",
            "Listing Control: List organic wheat, rice, pulses, and vegetables with photos and descriptions.",
            "Autonomous Pricing: Farmers set their own prices based on quality, bypassing forced broker rates."
        ],
        title_color=emerald
    )
    
    add_card(
        slide7,
        Inches(0.75) + 1 * (col_width + col_gap), col_top, col_width, col_height,
        "2. Consumer Purchase",
        [
            "Direct Ordering: Consumers and local retailers browse verified organic products.",
            "Distance Filters: Lists crops based on closeness to buyer to minimize freight costs.",
            "Zero Middleman Inflation: End-users get fresh, chemical-free food at direct prices."
        ],
        title_color=light_emerald
    )
    
    add_card(
        slide7,
        Inches(0.75) + 2 * (col_width + col_gap), col_top, col_width, col_height,
        "3. Quality Assurance",
        [
            "Traceability Badges: Products linked directly to soil test records receive premium 'Soil-Certified' tags.",
            "Consumer Reviews: Buyers rate produce quality, establishing seller reputation.",
            "Escrow Security: Payments held until buyer confirms receipt of produce."
        ],
        title_color=white
    )

    # -------------------------------------------------------------
    # SLIDE 8: Smart Logistics & Auto-Assigned Cheapest Transporters
    # -------------------------------------------------------------
    logistics_steps = [
        (
            "1. Order & Accept",
            [
                "User purchases organic crop or mill fertilizer.",
                "Seller accepts order, triggering a transport request."
            ],
            emerald
        ),
        (
            "2. Automated Routing",
            [
                "Logistics engine calculates route between seller and buyer.",
                "Determines distance, traffic grids, and transport load class."
            ],
            light_emerald
        ),
        (
            "3. Least-Cost Matching",
            [
                "System queries regional transporter rates on that specific route.",
                "Auto-assigns the job to the transporter bidding the lowest price."
            ],
            white
        ),
        (
            "4. Tracked Delivery",
            [
                "Transporter collects shipment using matched vehicle (motorcycle to truck).",
                "Updates GPS coordinate markers on Leaflet Maps in real-time."
            ],
            emerald
        )
    ]
    add_diagram_slide("Smart Logistics: Automated Least-Cost Route Matching", logistics_steps)

    # -------------------------------------------------------------
    # SLIDE 9: Detailed Revenue Model
    # -------------------------------------------------------------
    slide9 = add_blank_slide()
    add_header(slide9, "Ecosystem Revenue Model: Sustainable & Scalable")
    
    # 4 columns setup
    col_width_4 = Inches(2.7)
    col_gap_4 = Inches(0.34)
    col_top_4 = Inches(1.85)
    col_height_4 = Inches(4.8)
    
    add_card(
        slide9,
        Inches(0.75) + 0 * (col_width_4 + col_gap_4), col_top_4, col_width_4, col_height_4,
        "1. Consumer Trade Cut",
        [
            "Flat 3% Commission: Levied on successful direct-to-buyer transactions for organic crops.",
            "Paid by Buyer: Subsidizes platform operations without eating into farmer profits.",
            "Secure Escrow: Automates payouts via integrated payment gateways."
        ],
        title_color=emerald
    )
    
    add_card(
        slide9,
        Inches(0.75) + 1 * (col_width_4 + col_gap_4), col_top_4, col_width_4, col_height_4,
        "2. Mill Fertilizer Margin",
        [
            "Sourcing Margins: Split delivery margins on bulk industrial organic byproduct contracts.",
            "Volume Discounts: Sells at massive volumes to farmer cooperatives, creating high cash flow.",
            "Mill Referral Fees: Direct advertising margins paid by partnered mills."
        ],
        title_color=light_emerald
    )
    
    add_card(
        slide9,
        Inches(0.75) + 2 * (col_width_4 + col_gap_4), col_top_4, col_width_4, col_height_4,
        "3. Equipment Rental fee",
        [
            "5% P2P Service Fee: Charged on peer-to-peer machinery rentals.",
            "Platform Insurance: Part of the fee funds primary equipment damage micro-insurance.",
            "Transaction Value: Generates revenue from idle resources without inventory costs."
        ],
        title_color=white
    )
    
    add_card(
        slide9,
        Inches(0.75) + 3 * (col_width_4 + col_gap_4), col_top_4, col_width_4, col_height_4,
        "4. OEM Lead Fees",
        [
            "Brand Partnerships: Referral fees paid by tractor and machinery manufacturers (OEMs) for leads.",
            "Sponsored Listing Ads: Manufacturers pay to feature new machinery on the shopping portal.",
            "Bulk Buying Tenders: platform aggregates demand to secure lower OEM prices."
        ],
        title_color=gold
    )

    # -------------------------------------------------------------
    # SLIDE 10: Technical Stack Overview (Tech Page 1)
    # -------------------------------------------------------------
    slide10 = add_blank_slide()
    add_header(slide10, "Technology Stack & Core System Interfaces")
    
    add_card(
        slide10,
        Inches(0.75), Inches(1.85), Inches(5.6), Inches(4.8),
        "Frontend Interface (Client)",
        [
            "React.js (Vite): Optimized component-based UI offering fast rendering cycles.",
            "Interactive Leaflet GIS: Real-time map overlays displaying farm boundaries, testing labs, and live transporter routes.",
            "Tailored Visual CSS: High-contrast responsive design optimized for rural mobile web viewports.",
            "Localized Language Interceptor: Translates Hinglish queries (e.g. 'ganna beej') to standard keywords to aid navigation."
        ],
        title_color=emerald
    )
    
    add_card(
        slide10,
        Inches(6.98), Inches(1.85), Inches(5.6), Inches(4.8),
        "Backend & System Integrations",
        [
            "Node.js & Express.js: Event-driven API server handling high concurrent requests asynchronously.",
            "Passport.js Auth: Standardized session-based user authentication and security.",
            "Grok AI (X.AI API): Processes raw N-P-K & pH soil data using custom agronomic templates to generate suggestions.",
            "Cloudinary & Resend: Secure report hosting (PDFs) and direct email alerts/OTP codes."
        ],
        title_color=light_emerald
    )

    # -------------------------------------------------------------
    # SLIDE 11: Core Architecture & Data Syncing (Tech Page 2)
    # -------------------------------------------------------------
    slide11 = add_blank_slide()
    add_header(slide11, "Data Architecture, Caching & System Resilience")
    
    # 3 columns for tech execution
    col_width_t = Inches(3.6)
    col_gap_t = Inches(0.5)
    col_top_t = Inches(1.85)
    col_height_t = Inches(4.8)
    
    add_card(
        slide11,
        Inches(0.75) + 0 * (col_width_t + col_gap_t), col_top_t, col_width_t, col_height_t,
        "1. Database Layer",
        [
            "MongoDB Atlas: Cloud document store mapping highly nested schemas.",
            "Dedicated Collections:\n  - User (Roles, Pins)\n  - Listing (Produce, Rent)\n  - Order (Transporter details)\n  - SoilTest (Coordinates, Lab results, Grok advice)"
        ],
        title_color=emerald
    )
    
    add_card(
        slide11,
        Inches(0.75) + 1 * (col_width_t + col_gap_t), col_top_t, col_width_t, col_height_t,
        "2. Caching Engine",
        [
            "Redis (ioredis): Fast in-memory cache layer deployed locally in the cluster.",
            "OTP Expirations: Enforces short TTL limits for user registrations.",
            "Search Optimization: Caches heavy database queries for search terms."
        ],
        title_color=light_emerald
    )
    
    add_card(
        slide11,
        Inches(0.75) + 2 * (col_width_t + col_gap_t), col_top_t, col_width_t, col_height_t,
        "3. Local Resilience Layer",
        [
            "Local Memory Fallback: If Redis fails, API automatically switches to in-memory JS Map objects.",
            "Database Fallbacks: Launches local mongo-memory-server if Atlas becomes unreachable.",
            "Offline Operational Uptime: Guarantees zero transaction drops in remote areas."
        ],
        title_color=white
    )

    # -------------------------------------------------------------
    # SLIDE 12: Why Trishastik Bharat is the Best & Unique
    # -------------------------------------------------------------
    slide12 = add_blank_slide()
    add_header(slide12, "Why Trishastik Bharat is the Unique & Best Solution")
    
    # Left Column: Key Innovations
    add_card(
        slide12,
        Inches(0.75), Inches(1.85), Inches(5.6), Inches(4.8),
        "Core Project Achievements",
        [
            "True Disintermediation: Completely removes middleman margins, increasing farmer payouts by 25%+ directly.",
            "Direct Industrial Partnerships: Connecting mills directly with farms bypasses retail chemical chemical-only marketplaces.",
            "Enforced Public Accountability: Portal integration converts unmonitored free lab tests into audited, timely diagnostic runs.",
            "Auto-Optimized Freight Costs: Least-cost routing algorithm ensures transportation fees are kept to a bare minimum."
        ],
        title_color=white
    )

    # Right Column: Why We Win
    add_card(
        slide12,
        Inches(6.98), Inches(1.85), Inches(5.6), Inches(4.8),
        "Why Trishastik is Best in Class",
        [
            "Circular AgTech Economy: Unifies seeds, fertilizer, rentals, testing, sales, and logistics in a single cooperative loop.",
            "Asset-Light Operations: Uses peer-to-peer equipment sharing and local transporter networks, requiring zero vehicle capital.",
            "Science-Driven Agriculture: Combines official laboratory tests with Grok AI language diagnostics, making data readable in local terms.",
            "Built for Rural Scaling: Lightweight MERN stack works seamlessly on low-bandwidth rural connections."
        ],
        title_color=emerald
    )

    # Save presentation
    output_filename = "Trishastik_Bharat_Presentation.pptx"
    saved = False
    for suffix in ["", "_v2", "_v3", "_v4", "_v5", "_final"]:
        name = f"Trishastik_Bharat_Presentation{suffix}.pptx" if suffix else "Trishastik_Bharat_Presentation.pptx"
        try:
            prs.save(name)
            output_filename = name
            saved = True
            print(f"Presentation saved successfully as '{output_filename}'")
            break
        except PermissionError:
            continue
    if not saved:
        import time
        output_filename = f"Trishastik_Bharat_Presentation_{int(time.time())}.pptx"
        prs.save(output_filename)
        print(f"All standard names locked. Saved successfully as '{output_filename}'")
    
    # Also save to the extr directory explicitly to ensure it gets updated there
    try:
        extr_path = os.path.join("extr", "Trishastik_Bharat_Presentation.pptx")
        prs.save(extr_path)
        print(f"Presentation copy saved in extr directory: {extr_path}")
    except Exception as e:
        print(f"Failed to copy to extr: {e}")

if __name__ == "__main__":
    create_presentation()
